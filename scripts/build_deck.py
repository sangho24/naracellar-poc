"""
매니저 review용 4장 deck 생성 v2 — Why AX Node 중심.
디자인 최소 (사용자가 별도 폴리시). 위치/컨텐츠만 정확히.

산출물: output/manager_review_deck_v2.pptx
실행: py scripts/build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Canvas (16:9) ───────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

# ─── Color palette (디자인 폴리시 전 안전한 톤) ───────────────────────────────
INK_900   = RGBColor(0x1A, 0x14, 0x10)
INK_700   = RGBColor(0x3D, 0x34, 0x2D)
INK_500   = RGBColor(0x6F, 0x67, 0x60)
INK_300   = RGBColor(0xA8, 0xA2, 0x9E)
LINE      = RGBColor(0xE8, 0xE5, 0xDE)
CANVAS    = RGBColor(0xFA, 0xF9, 0xF6)
BURGUNDY  = RGBColor(0x9A, 0x25, 0x3C)
GOLD      = RGBColor(0xB8, 0x9A, 0x65)
DANGER    = RGBColor(0xC7, 0x37, 0x3B)
POSITIVE  = RGBColor(0x2D, 0x7D, 0x3F)
WARNING   = RGBColor(0xB4, 0x53, 0x09)
SOFT_BG   = RGBColor(0xFB, 0xEE, 0xF1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

KO_FONT = "맑은 고딕"


def add_text(slide, left, top, width, height, text,
             size=12, bold=False, color=INK_900, align=PP_ALIGN.LEFT,
             font_name=KO_FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, lines,
                size=11, color=INK_700, bullet="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = f"{bullet} {ln}"
        run.font.size = Pt(size)
        run.font.name = KO_FONT
        run.font.color.rgb = color
    return tb


def add_box(slide, left, top, width, height, fill=None, line=LINE, line_w=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.05
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_w)
    shape.shadow.inherit = False
    return shape


def add_header(slide, eyebrow, title, subtitle, page_num):
    """모든 슬라이드 상단 공통 헤더."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(0.08), Inches(0.4))
    bar.fill.solid(); bar.fill.fore_color.rgb = BURGUNDY
    bar.line.fill.background()

    add_text(slide, Inches(0.7), Inches(0.42), Inches(3), Inches(0.3),
             eyebrow.upper(), size=10, bold=True, color=BURGUNDY)
    add_text(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.7),
             title, size=26, bold=True, color=INK_900)
    add_text(slide, Inches(0.5), Inches(1.45), Inches(12.3), Inches(0.5),
             subtitle, size=12, color=INK_500)

    sep = slide.shapes.add_connector(1, Inches(0.5), Inches(2.0), Inches(12.8), Inches(2.0))
    sep.line.color.rgb = LINE
    sep.line.width = Pt(0.5)

    add_text(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
             "삼일회계법인 AX Node", size=9, color=INK_500)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{page_num} / 4", size=9, color=INK_500, align=PP_ALIGN.RIGHT)


def add_blank_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


# ============================================================================
# SLIDE 1 — As-Is: 재무 위기 × 시장 구조 변화
# ============================================================================
s = add_blank_slide()
add_header(s, "As-Is",
           "재무·시장 양쪽에서 동시에 압박받는 구조",
           "현금흐름 위기 + 한국 와인 시장 구조 변화가 나라셀라 강점 카테고리를 직격",
           1)

add_text(s, Inches(0.5), Inches(2.2), Inches(6), Inches(0.4),
         "재무적 압박", size=14, bold=True, color=BURGUNDY)
add_bullets(s, Inches(0.5), Inches(2.65), Inches(6), Inches(4.0), [
    "부채 316억 vs 현금 44억 — 7배 차이 · 이자 연 16억",
    "2024 영업손실 -34억 (annual, 적자 전환) · 2025.3Q -2.3억",
    "시총 1,288억 → 310억 (-76%) · IPO 공모가 -87%",
    "재고보유일수 348일 (2022 215일 → +133일)",
    "    - 미착상품 102억 (전체 재고의 25%)",
    "    - 재고평가충당금 2.86억",
    "외화부채 73억 (5% 변동 시 3.6억 손익, 수동 헤지)",
    "매출채권 25.4억 (연체율 18.6%) · 대손 매년 +2억",
    "자체 IT 시도 실패 — 1KMWINE 와인원 45억(목표 100억),",
    "    물류센터 무기한 보류",
], size=11, color=INK_700)

add_text(s, Inches(6.9), Inches(2.2), Inches(6), Inches(0.4),
         "한국 와인 시장 구조 변화", size=14, bold=True, color=BURGUNDY)
add_bullets(s, Inches(6.9), Inches(2.65), Inches(6), Inches(3.0), [
    "수입량 -20.4% (2021 76,575t → 2024 56,542t)",
    "카테고리 변화: 레드 69%→51%, 화이트 19%→34%",
    "양극화 — 초고가 견조, 중간(3~10만) 급감,",
    "    가성비는 편의점 흡수",
    "편의점 와인 폭발 (디아블로 단일 200만병)",
    "MZ 이탈: CU 하이볼 8.3%→38.6% (2022→2024)",
    "    무알코올 +9%/yr (2030년까지)",
], size=11, color=INK_700)

add_box(s, Inches(6.9), Inches(4.5), Inches(6), Inches(1.45), fill=CANVAS, line=LINE)
add_text(s, Inches(7.1), Inches(4.55), Inches(5.6), Inches(0.3),
         "경쟁사 비교 (2024 영업이익)", size=10, bold=True, color=INK_500)
add_bullets(s, Inches(7.1), Inches(4.8), Inches(5.6), Inches(1.1), [
    "신세계L&B: -52억 (적자 전환)   · 금양: -82.5%",
    "아영FBC: +40% (직영매장·고마진 — 유일한 영업익 ↑)",
    "나라셀라: -34억 (적자 전환)",
], size=10, color=INK_700)

add_box(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.85),
        fill=SOFT_BG, line=BURGUNDY)
add_text(s, Inches(0.8), Inches(6.2), Inches(11.8), Inches(0.3),
         "나라셀라 강점(레드 mid-premium)이 가장 약한 위치에 놓임",
         size=12, bold=True, color=BURGUNDY)
add_text(s, Inches(0.8), Inches(6.5), Inches(11.8), Inches(0.4),
         "외부 트렌드는 통제 불가 → 내부 의사결정 가시성·자동화·AI 매칭으로 turnaround",
         size=10.5, color=INK_700)


# ============================================================================
# SLIDE 2 — To-Be: AI Agent로 1,955 × 570 × 10+ 매칭
# ============================================================================
s = add_blank_slide()
add_header(s, "To-Be",
           "포트폴리오 전체가 살아 있는 회사",
           "AI Agent의 영역 — 영업사원의 머리로 풀 수 없는 조합 최적화",
           2)

add_text(s, Inches(0.5), Inches(2.15), Inches(12.3), Inches(0.6),
         "1,955 SKU × 570 거래처 × 10+ 채널",
         size=24, bold=True, color=BURGUNDY, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(2.75), Inches(12.3), Inches(0.4),
         "이 조합을 영업사원의 머리로 풀 수는 없습니다. 이건 AI Agent의 영역입니다.",
         size=12, color=INK_700, align=PP_ALIGN.CENTER)

EFFECTS = [
    ("Long Tail 소진",     "비인기 SKU가 맞는 거래처에 연결\n→ 재고보유일수 정상화"),
    ("거래처 차별화",       "\"우리 매장에 맞는 와인을\n추천받을 수 있다\" 인식 확대"),
    ("얼로케이션 선순환",   "\"한국에서 잘 팔린다\" → 와이너리\n물량 확대 → 포트폴리오 강화"),
    ("상장 비전 달성",      "인기 와인 + Long Tail\n→ 매출 2,500억 경로 재활성화"),
]
card_top = Inches(3.45)
card_w = Inches(2.9)
card_h = Inches(1.6)
gap = Inches(0.15)
total_w = card_w * 4 + gap * 3
start_x = (SW - total_w) // 2

for i, (title, body) in enumerate(EFFECTS):
    left = start_x + (card_w + gap) * i
    add_box(s, left, card_top, card_w, card_h, fill=CANVAS, line=LINE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, card_top, card_w, Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = BURGUNDY
    bar.line.fill.background()
    add_text(s, left + Inches(0.2), card_top + Inches(0.15), Inches(0.4), Inches(0.3),
             f"#{i+1}", size=10, bold=True, color=GOLD)
    add_text(s, left + Inches(0.2), card_top + Inches(0.4), card_w - Inches(0.4), Inches(0.4),
             title, size=13, bold=True, color=INK_900)
    add_text(s, left + Inches(0.2), card_top + Inches(0.85), card_w - Inches(0.4), Inches(0.7),
             body, size=10, color=INK_700)

add_text(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.4),
         "왜 이게 정량적으로 가능한가 — 자체 리서치 (PoC + 글로벌 트렌드)",
         size=12, bold=True, color=INK_900)

add_box(s, Inches(0.5), Inches(5.7), Inches(6), Inches(1.3), fill=None, line=LINE)
add_text(s, Inches(0.65), Inches(5.8), Inches(5.7), Inches(0.3),
         "740 SKU 외부 데이터 분석 (국가별 SKU vs 매출 gap)",
         size=10, bold=True, color=BURGUNDY)
add_bullets(s, Inches(0.65), Inches(6.1), Inches(5.7), Inches(0.85), [
    "프랑스 +14.8%p — Long Tail 위험 #1",
    "이탈리아 +8.4%p — #2 · 미국 -11%p — under-curated",
    "큐레이션 mismatch가 정량으로 드러남",
], size=9, color=INK_700)

add_box(s, Inches(6.83), Inches(5.7), Inches(6), Inches(1.3), fill=None, line=LINE)
add_text(s, Inches(6.98), Inches(5.8), Inches(5.7), Inches(0.3),
         "AI 도입 5단 효과 (단순 생산성 X)",
         size=10, bold=True, color=BURGUNDY)
add_bullets(s, Inches(6.98), Inches(6.1), Inches(5.7), Inches(0.85), [
    "Productivity · Decision Quality · Revenue Growth",
    "Risk Mitigation · Strategic Optionality (M&A·얼로케이션)",
    "글로벌 트렌드 정합 — AI 마케팅·1차 데이터·하이퍼 개인화",
], size=9, color=INK_700)


# ============================================================================
# SLIDE 3 — Roadmap (이 용역의 업무 Scope)
# ============================================================================
s = add_blank_slide()
add_header(s, "Scope",
           "5.5개월, 3 Phase — 운영 시스템까지 인계",
           "거창한 DT 청사진이 아니라, 본 용역에서 우리가 만들어 인계할 시스템·산출물 단위",
           3)

PHASES = [
    {
        "num": "Phase 1", "weeks": "0~6주", "title": "PoC 본격 도입",
        "items": [
            "데이터 정제 — 1,613 active SKU × 570 거래처",
            "Knowledge Graph 스키마·시드 (Neo4j)",
            "이상 감지 + Alert 자동화",
            "Root Cause AI Agent (LangGraph 4 노드)",
            "운영 Dashboard 환경 이전",
        ],
        "deliv": "운영 시스템 (alpha) + 데이터 정제 매뉴얼",
    },
    {
        "num": "Phase 2", "weeks": "6~14주", "title": "핵심 자동화",
        "items": [
            "수요예측 모델 (SKU × 채널)",
            "거래처 신용 스코링 모델",
            "환리스크 자동 추적 + 헤지 알림",
            "What-If 시뮬레이터",
            "재무 보고서 자동화 (월/분기)",
        ],
        "deliv": "자동화 모듈 + 운영 Runbook + 사용자 교육 1회",
    },
    {
        "num": "Phase 3", "weeks": "14~22주", "title": "AI 추천 Agent",
        "items": [
            "거래처별 맞춤 와인 추천 모델",
            "영업사원 미팅 브리핑 자동 생성",
            "신규 SKU 론칭 최적화 시뮬레이션",
            "와이너리 얼로케이션 협상 데이터",
        ],
        "deliv": "AI 추천 시스템 + 영업 매뉴얼 + 12개월 운영 지원",
    },
]

phase_top = Inches(2.2)
phase_w = Inches(4.0)
phase_h = Inches(3.6)
phase_gap = Inches(0.15)
phase_total = phase_w * 3 + phase_gap * 2
phase_start = (SW - phase_total) // 2

for i, p in enumerate(PHASES):
    left = phase_start + (phase_w + phase_gap) * i
    add_box(s, left, phase_top, phase_w, phase_h, fill=CANVAS, line=LINE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, phase_top, phase_w, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = BURGUNDY
    bar.line.fill.background()
    add_text(s, left + Inches(0.25), phase_top + Inches(0.2), Inches(2), Inches(0.3),
             p["num"], size=11, bold=True, color=BURGUNDY)
    add_text(s, left + Inches(0.25), phase_top + Inches(0.5), Inches(3.5), Inches(0.3),
             p["weeks"], size=9, color=INK_500)
    add_text(s, left + Inches(0.25), phase_top + Inches(0.8), Inches(3.5), Inches(0.4),
             p["title"], size=14, bold=True, color=INK_900)
    add_bullets(s, left + Inches(0.25), phase_top + Inches(1.3), Inches(3.5), Inches(1.8),
                p["items"], size=9.5, color=INK_700)
    deliv_top = phase_top + Inches(3.0)
    add_text(s, left + Inches(0.25), deliv_top, Inches(3.5), Inches(0.2),
             "산출물", size=8.5, bold=True, color=GOLD)
    add_text(s, left + Inches(0.25), deliv_top + Inches(0.2), Inches(3.5), Inches(0.4),
             p["deliv"], size=9, color=INK_700)

add_box(s, Inches(0.5), Inches(6.05), Inches(8), Inches(0.95), fill=None, line=LINE)
add_text(s, Inches(0.7), Inches(6.15), Inches(7.6), Inches(0.3),
         "인력·체계", size=10, bold=True, color=BURGUNDY)
add_text(s, Inches(0.7), Inches(6.4), Inches(7.6), Inches(0.5),
         "데이터 엔지니어 2 · AI 엔지니어 2 · 도메인 컨설턴트 1 · PM 1\n격주 운영 위원회 · 월 1회 회장님 보고 (재무·운영 통합)",
         size=9.5, color=INK_700)

add_box(s, Inches(8.7), Inches(6.05), Inches(4.13), Inches(0.95),
        fill=SOFT_BG, line=BURGUNDY)
add_text(s, Inches(8.9), Inches(6.2), Inches(3.8), Inches(0.6),
         "Phase 종료 시점마다\n운영 가능한 시스템 인계",
         size=11, bold=True, color=BURGUNDY)


# ============================================================================
# SLIDE 4 — Why 삼일 AX Node (가장 중요)
# ============================================================================
s = add_blank_slide()
add_header(s, "Why AX Node",
           "디지털 전환을 회계법인이 — 왜 우리만 가능한가",
           "같은 PoC를 다른 컨설팅펌도 만들 수 있지만, 회장님 언어로 IR을 데이터로 증명할 수 있는 곳은 삼일 AX Node뿐",
           4)

why_cards = [
    {
        "tag": "#1 — MAIN",
        "title": "재무 언어로 소통 가능한 유일 파트너",
        "body_lines": [
            "마승철 회장 = 디아지오 본사·코리아 CFO 출신",
            "사외이사 최범수 = 前 KB은행 부행장",
            "회장님 IR 발언과 우리 효과 항목 1:1 매칭 (아래 표)",
            "IT 펌이 못 하는 것 — CFO 언어로 IR을 데이터로 증명",
        ],
        "color": BURGUNDY,
    },
    {
        "tag": "#2",
        "title": "사내 Sponsorship + 글로벌 PwC Wineries Network",
        "body_lines": [
            "이미 사내 제휴 등록 (2026 4~5월) — Trust 단계 통과",
            "Constellation·Treasury 등 글로벌 와인 그룹 audit 경험",
            "와이너리 얼로케이션 — 미·프·이 PwC member firms",
            "글로벌 distributor KPI·운영 벤치마크 접근",
        ],
        "color": GOLD,
    },
    {
        "tag": "#3",
        "title": "Audit × AX 통합 + 외부 파트너 수용 명분",
        "body_lines": [
            "KOSDAQ 405920 상장사 — audit + AX 통합 (빅4 최강)",
            "1KMWINE 자체 IT 실패 → 외부 파트너 수용 명분 확보",
            "단순 SI가 아닌 재무 일체형 솔루션",
            "아영FBC는 H/W로, 우리는 데이터·AI S/W로 같은 곳에",
        ],
        "color": INK_700,
    },
]

why_top = Inches(2.15)
why_w = Inches(4.0)
why_h = Inches(2.4)
why_gap = Inches(0.15)
why_total = why_w * 3 + why_gap * 2
why_start = (SW - why_total) // 2

for i, c in enumerate(why_cards):
    left = why_start + (why_w + why_gap) * i
    add_box(s, left, why_top, why_w, why_h, fill=CANVAS, line=LINE)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, why_top, why_w, Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = c["color"]
    bar.line.fill.background()
    add_text(s, left + Inches(0.25), why_top + Inches(0.2), Inches(3.5), Inches(0.3),
             c["tag"], size=10, bold=True, color=c["color"])
    add_text(s, left + Inches(0.25), why_top + Inches(0.5), Inches(3.5), Inches(0.6),
             c["title"], size=12, bold=True, color=INK_900)
    add_bullets(s, left + Inches(0.25), why_top + Inches(1.15), Inches(3.5), Inches(1.2),
                c["body_lines"], size=9.5, color=INK_700)

# 매칭 표
table_top = Inches(4.75)
add_text(s, Inches(0.5), table_top, Inches(12.3), Inches(0.3),
         "회장님 IR 발언 ↔ 본 용역 효과 1:1 매칭",
         size=11, bold=True, color=BURGUNDY)

table_rows = [
    ("회장님 IR 발언 (2025-03)", "본 용역의 직접 효과"),
    ("\"발주 예산제·재고 선순환\"", "수요예측 모델 → 재고 348→250일 (100억 회수)"),
    ("\"판관비 개선\"", "IT 통합·자동화 → 운영비 ~30% 절감"),
    ("\"영업이익 턴어라운드\"", "비용절감 합계 연 7~10억 → 흑자 전환"),
    ("\"포트폴리오 균형\"", "AI 추천 → Long Tail 매출 +5~10%"),
]
table = s.shapes.add_table(rows=5, cols=2,
                            left=Inches(0.5), top=Inches(5.1),
                            width=Inches(12.3), height=Inches(1.7)).table
table.columns[0].width = Inches(5.5)
table.columns[1].width = Inches(6.8)
for r, (l, rval) in enumerate(table_rows):
    for c, val in enumerate((l, rval)):
        cell = table.cell(r, c)
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = val
        run.font.name = KO_FONT
        run.font.size = Pt(9.5 if r > 0 else 10)
        run.font.bold = (r == 0)
        if r == 0:
            run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = BURGUNDY
        else:
            run.font.color.rgb = INK_900 if c == 1 else INK_700
            cell.fill.solid()
            cell.fill.fore_color.rgb = CANVAS if r % 2 == 1 else WHITE

add_box(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.6),
        fill=BURGUNDY, line=BURGUNDY)
add_text(s, Inches(0.7), Inches(6.92), Inches(11.9), Inches(0.5),
         "회장님이 IR에서 말씀하시는 모든 키워드 — 발주 예산제·재고 선순환·판관비·턴어라운드 — 를 데이터로 증명할 수 있는 곳은 삼일 AX Node뿐.",
         size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─── Save ────────────────────────────────────────────────────────────────────
import os
OUT_DIR = os.path.join(os.path.dirname(__file__), "..", "output")
# v2 파일명 사용 — 기존 pptx가 PowerPoint에 열려 있어도 충돌 안 함
OUT = os.path.join(OUT_DIR, "manager_review_deck_v2.pptx")
prs.save(OUT)
print(f"[OK] saved → {os.path.abspath(OUT)}")
